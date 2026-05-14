"""
PVM URL Admin — PPTX export
Converts the 5 admin views into a 6-slide deck (cover + 5 screens).
Uses footer-rail + cursor-flow discipline from pptx-html-fidelity-audit skill.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Canvas (16:9) ──────────────────────────────────────────────────────────────
CANVAS_W      = Inches(13.333)
CANVAS_H      = Inches(7.5)

# ── Margins ────────────────────────────────────────────────────────────────────
MARGIN_X      = Inches(0.55)
MARGIN_TOP    = Inches(0.55)
CONTENT_LEFT  = MARGIN_X
CONTENT_W     = CANVAS_W - MARGIN_X * 2

# ── Vertical rails ─────────────────────────────────────────────────────────────
CONTENT_MAX_Y = Inches(6.70)   # nothing in content area may cross this
FOOTER_TOP    = Inches(6.85)   # footer row pinned here
FOOTER_H      = Inches(0.22)
CHROME_TOP    = Inches(0.20)
CHROME_H      = Inches(0.18)

# ── PVM brand palette ──────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x0D, 0x1F, 0x35)
C_ACCENT  = RGBColor(0xF5, 0xC4, 0x00)
C_BG      = RGBColor(0xEE, 0xF1, 0xF5)
C_SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
C_FG      = RGBColor(0x0D, 0x1F, 0x35)
C_MUTED   = RGBColor(0x6B, 0x72, 0x80)
C_BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
C_TEAL    = RGBColor(0x02, 0x84, 0xC7)
C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_NAVY_60 = RGBColor(0x68, 0x77, 0x8E)

# ── Fonts ──────────────────────────────────────────────────────────────────────
FONT_SANS = "Segoe UI"
FONT_MONO = "Cascadia Code"


# ══════════════════════════════════════════════════════════════════════════════
#  Cursor
# ══════════════════════════════════════════════════════════════════════════════

class Cursor:
    def __init__(self, y_start=None, cap=CONTENT_MAX_Y):
        self.y = y_start if y_start is not None else MARGIN_TOP
        self.cap = cap

    def take(self, h, gap=Inches(0.12), label=""):
        top = self.y
        self.y = top + h + gap
        if self.y > self.cap:
            raise OverflowError(
                f"Cursor exceeded rail at '{label}': "
                f"bottom={top+h:.4f} rail={self.cap:.4f}"
            )
        return top


def hero_layout(blocks):
    total = sum(h + g for h, g in blocks)
    y_start = (CANVAS_H - total) / 2
    return Cursor(y_start=y_start, cap=CANVAS_H - FOOTER_H - Inches(0.15))


# ══════════════════════════════════════════════════════════════════════════════
#  Low-level helpers
# ══════════════════════════════════════════════════════════════════════════════

def _set_run_font(run, font_name, size_pt, bold=False, italic=False,
                  color=None, all_caps=False, letter_spacing_pt=None):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if all_caps:
        run.font.all_caps = True
    # letter spacing via rPr spacing attribute (in hundredths of a point)
    if letter_spacing_pt is not None:
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", str(int(letter_spacing_pt * 100)))
    # Ensure latin typeface slot is set explicitly to avoid Calibri fallback
    rpr = run._r.get_or_add_rPr()
    latin = rpr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rpr, qn("a:latin"))
    latin.set("typeface", font_name)
    # ea slot — set to +mj-ea or a safe fallback to avoid JhengHei
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", "+mn-ea")


def add_textbox(slide, left, top, width, height, text,
                font=FONT_SANS, size_pt=11, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, all_caps=False,
                letter_spacing_pt=None, word_wrap=True,
                line_spacing_pt=None, name=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    if name:
        txb.name = name
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing_pt:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, _qn("a:lnSpc"))
        spcPts = etree.SubElement(lnSpc, _qn("a:spcPts"))
        spcPts.set("val", str(int(line_spacing_pt * 100)))
    run = p.add_run()
    run.text = text
    _set_run_font(run, font, size_pt, bold=bold, italic=italic,
                  color=color, all_caps=all_caps,
                  letter_spacing_pt=letter_spacing_pt)
    return txb


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width_pt=0.5, name=None):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if name:
        shape.name = name
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = _Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def paint_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, left_text, right_text, theme="light"):
    color = C_NAVY_60 if theme == "light" else RGBColor(0x9B, 0xA0, 0xA6)
    add_textbox(slide, CONTENT_LEFT, FOOTER_TOP,
                CONTENT_W / 2, FOOTER_H, left_text,
                font=FONT_MONO, size_pt=8, color=color,
                align=PP_ALIGN.LEFT, name="footer-left",
                letter_spacing_pt=1.0)
    add_textbox(slide, CONTENT_LEFT + CONTENT_W / 2, FOOTER_TOP,
                CONTENT_W / 2, FOOTER_H, right_text,
                font=FONT_MONO, size_pt=8, color=color,
                align=PP_ALIGN.RIGHT, name="footer-right",
                letter_spacing_pt=1.0)


def add_chrome(slide, left_text, right_text, theme="light"):
    color = C_MUTED
    add_textbox(slide, CONTENT_LEFT, CHROME_TOP,
                CONTENT_W / 2, CHROME_H, left_text,
                font=FONT_MONO, size_pt=7.5, color=color,
                align=PP_ALIGN.LEFT, name="chrome-left",
                all_caps=True, letter_spacing_pt=1.5)
    add_textbox(slide, CONTENT_LEFT + CONTENT_W / 2, CHROME_TOP,
                CONTENT_W / 2, CHROME_H, right_text,
                font=FONT_MONO, size_pt=7.5, color=color,
                align=PP_ALIGN.RIGHT, name="chrome-right",
                all_caps=True, letter_spacing_pt=1.5)


# ══════════════════════════════════════════════════════════════════════════════
#  Shared sidebar mock
# ══════════════════════════════════════════════════════════════════════════════

SIDEBAR_W = Inches(1.9)

def add_sidebar(slide, active_view="dashboard"):
    # Navy sidebar background (named chrome-* so the rail verifier exempts it)
    add_rect(slide, Inches(0), Inches(0), SIDEBAR_W, CANVAS_H,
             fill_color=C_NAVY, name="chrome-sidebar-bg")
    # Accent stripe at top
    add_rect(slide, Inches(0), Inches(0), SIDEBAR_W, Inches(0.04),
             fill_color=C_ACCENT, name="sidebar-accent-stripe")
    # Brand logo box
    add_rect(slide, Inches(0.18), Inches(0.14), Inches(0.32), Inches(0.32),
             fill_color=C_ACCENT, name="sidebar-logo-bg")
    add_textbox(slide, Inches(0.18), Inches(0.14), Inches(0.32), Inches(0.32),
                "P", font=FONT_SANS, size_pt=13, bold=True,
                color=C_NAVY, align=PP_ALIGN.CENTER, name="sidebar-logo")
    # Brand name
    add_textbox(slide, Inches(0.56), Inches(0.16), Inches(1.22), Inches(0.16),
                "PVM URL Admin", font=FONT_SANS, size_pt=8, bold=True,
                color=RGBColor(0xFF,0xFF,0xFF), name="sidebar-name")
    add_textbox(slide, Inches(0.56), Inches(0.33), Inches(1.22), Inches(0.12),
                "MANAGED REDIRECTS", font=FONT_SANS, size_pt=6,
                color=RGBColor(0x9B,0xA0,0xA6), letter_spacing_pt=1.2,
                name="sidebar-sub")

    # Nav items
    nav_items = [
        ("OVERVIEW", None, False),
        ("Dashboard", "dashboard", True),
        ("REDIRECTS", None, False),
        ("All Redirects", "redirects", True),
        ("New Redirect", "create", True),
        ("ORGANISATION", None, False),
        ("Tags & Categories", "tags", True),
        ("SYSTEM", None, False),
        ("Settings", None, True),
    ]

    y = Inches(0.62)
    section_gap = Inches(0.06)
    item_h = Inches(0.28)
    section_h = Inches(0.18)

    for label, view_id, is_item in nav_items:
        if not is_item:
            # Section header
            add_textbox(slide, Inches(0.18), y, Inches(1.58), section_h,
                        label, font=FONT_SANS, size_pt=6.5, bold=True,
                        color=RGBColor(0x72,0x80,0x90),
                        all_caps=True, letter_spacing_pt=1.5,
                        name=f"sidebar-section-{label.lower()}")
            y += section_h + Inches(0.02)
        else:
            is_active = view_id == active_view
            if is_active:
                add_rect(slide, Inches(0.10), y - Inches(0.02),
                         SIDEBAR_W - Inches(0.20), item_h,
                         fill_color=RGBColor(0x1A,0x2E,0x48),
                         name=f"sidebar-active-bg-{view_id}")
            txt_color = C_ACCENT if is_active else RGBColor(0xA0,0xAA,0xB8)
            add_textbox(slide, Inches(0.22), y, Inches(1.50), item_h,
                        label, font=FONT_SANS, size_pt=9,
                        bold=is_active, color=txt_color,
                        name=f"sidebar-item-{label.replace(' ','-').lower()}")
            y += item_h + Inches(0.02)

    # User row at bottom
    add_rect(slide, Inches(0.10), Inches(6.90), SIDEBAR_W - Inches(0.20),
             Inches(0.34), fill_color=RGBColor(0x1A,0x2C,0x40),
             name="sidebar-user-bg")
    add_rect(slide, Inches(0.18), Inches(6.94), Inches(0.26), Inches(0.26),
             fill_color=RGBColor(0x37,0x41,0x51), name="sidebar-av-bg")
    add_textbox(slide, Inches(0.18), Inches(6.94), Inches(0.26), Inches(0.26),
                "M", font=FONT_SANS, size_pt=9, bold=True,
                color=C_SURFACE, align=PP_ALIGN.CENTER, name="sidebar-av")
    add_textbox(slide, Inches(0.50), Inches(6.95), Inches(1.20), Inches(0.14),
                "Marius", font=FONT_SANS, size_pt=8, bold=True,
                color=RGBColor(0xCC,0xD0,0xD6), name="sidebar-user-name")
    add_textbox(slide, Inches(0.50), Inches(7.10), Inches(1.20), Inches(0.12),
                "Administrator", font=FONT_SANS, size_pt=7,
                color=RGBColor(0x72,0x80,0x90), name="sidebar-user-role")


# ══════════════════════════════════════════════════════════════════════════════
#  Shared topbar
# ══════════════════════════════════════════════════════════════════════════════

def add_topbar(slide, breadcrumb_active, slide_num, total):
    content_start = SIDEBAR_W
    topbar_h = Inches(0.48)
    add_rect(slide, content_start, Inches(0), CANVAS_W - content_start,
             topbar_h, fill_color=C_SURFACE, line_color=C_BORDER,
             line_width_pt=0.5, name="topbar-bg")
    add_textbox(slide, content_start + Inches(0.22), Inches(0.12),
                Inches(4), Inches(0.22),
                breadcrumb_active, font=FONT_SANS, size_pt=9, bold=True,
                color=C_FG, name="topbar-bc")
    # "New redirect" button stub
    btn_w = Inches(1.18)
    add_rect(slide, CANVAS_W - Inches(0.22) - btn_w, Inches(0.10),
             btn_w, Inches(0.28), fill_color=C_NAVY, name="topbar-btn-bg")
    add_textbox(slide, CANVAS_W - Inches(0.22) - btn_w, Inches(0.10),
                btn_w, Inches(0.28), "+ New redirect",
                font=FONT_SANS, size_pt=8, bold=True,
                color=C_SURFACE, align=PP_ALIGN.CENTER, name="topbar-btn")


# ══════════════════════════════════════════════════════════════════════════════
#  Shared page header
# ══════════════════════════════════════════════════════════════════════════════

def add_page_header(slide, eyebrow, title, desc, cursor):
    CONTENT_START = SIDEBAR_W + Inches(0.28)
    W = CANVAS_W - CONTENT_START - Inches(0.28)
    cursor.take(Inches(0.12), gap=0, label="topbar-gap")  # space below topbar
    t = cursor.take(Inches(0.14), gap=Inches(0.04), label="eyebrow")
    add_textbox(slide, CONTENT_START, t, W, Inches(0.14),
                eyebrow, font=FONT_SANS, size_pt=7.5, bold=True,
                color=C_MUTED, all_caps=True, letter_spacing_pt=1.2,
                name="pg-eyebrow")
    t = cursor.take(Inches(0.30), gap=Inches(0.04), label="title")
    add_textbox(slide, CONTENT_START, t, W, Inches(0.30),
                title, font=FONT_SANS, size_pt=22, bold=True,
                color=C_FG, name="pg-title")
    t = cursor.take(Inches(0.18), gap=Inches(0.18), label="desc")
    add_textbox(slide, CONTENT_START, t, W, Inches(0.18),
                desc, font=FONT_SANS, size_pt=10, color=C_MUTED,
                name="pg-desc")


# ══════════════════════════════════════════════════════════════════════════════
#  Badge helper
# ══════════════════════════════════════════════════════════════════════════════

BADGE_COLORS = {
    "Active":     (RGBColor(0xF0,0xFD,0xF4), RGBColor(0x15,0x80,0x3D)),
    "Expired":    (RGBColor(0xFE,0xF2,0xF2), RGBColor(0xB9,0x1C,0x1C)),
    "Paused":     (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
    "Campaign":   (RGBColor(0xFF,0xFB,0xEB), RGBColor(0xB4,0x53,0x09)),
    "Print / QR": (RGBColor(0xEF,0xF6,0xFF), RGBColor(0x1D,0x4E,0xD8)),
    "Referral":   (RGBColor(0xF0,0xFD,0xF4), RGBColor(0x15,0x80,0x3D)),
    "Event":      (RGBColor(0xF5,0xF3,0xFF), RGBColor(0x6D,0x28,0xD9)),
    "General":    (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
    "Products":   (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
    "Events":     (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
    "Partners":   (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
    "Campaigns":  (RGBColor(0xF3,0xF4,0xF6), RGBColor(0x4B,0x55,0x63)),
}

def add_badge(slide, left, top, label, name=None):
    w = Inches(0.72)
    h = Inches(0.18)
    bg, fg = BADGE_COLORS.get(label, (C_BG, C_MUTED))
    add_rect(slide, left, top, w, h, fill_color=bg, line_color=None,
             name=name or f"badge-{label}")
    add_textbox(slide, left, top + Inches(0.01), w, h - Inches(0.01),
                label, font=FONT_SANS, size_pt=7.5, bold=True,
                color=fg, align=PP_ALIGN.CENTER,
                name=(name or f"badge-{label}") + "-text")


# ══════════════════════════════════════════════════════════════════════════════
#  KPI card
# ══════════════════════════════════════════════════════════════════════════════

def add_kpi_card(slide, left, top, w, h, label, value, delta="", delta_color=None):
    add_rect(slide, left, top, w, h, fill_color=C_SURFACE,
             line_color=C_BORDER, line_width_pt=0.5, name=f"kpi-{label}-bg")
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), w - Inches(0.28),
                Inches(0.14), label.upper(),
                font=FONT_SANS, size_pt=7, bold=True, color=C_MUTED,
                all_caps=True, letter_spacing_pt=0.8, name=f"kpi-{label}-lbl")
    add_textbox(slide, left + Inches(0.14), top + Inches(0.30), w - Inches(0.28),
                Inches(0.32), value,
                font=FONT_SANS, size_pt=26, bold=True, color=C_FG,
                name=f"kpi-{label}-val")
    if delta:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.66), w - Inches(0.28),
                    Inches(0.16), delta,
                    font=FONT_SANS, size_pt=8, color=delta_color or C_MUTED,
                    name=f"kpi-{label}-delta")


# ══════════════════════════════════════════════════════════════════════════════
#  Section card helper
# ══════════════════════════════════════════════════════════════════════════════

def add_card_bg(slide, left, top, width, height, name=None):
    return add_rect(slide, left, top, width, height,
                    fill_color=C_SURFACE, line_color=C_BORDER,
                    line_width_pt=0.5, name=name or "card-bg")

def add_card_header(slide, left, top, width, title, subtitle=None, name=None):
    h = Inches(0.36) if subtitle else Inches(0.28)
    add_rect(slide, left, top, width, h,
             fill_color=C_SURFACE, line_color=C_BORDER,
             line_width_pt=0.5, name=(name or "card-hd") + "-bg")
    add_textbox(slide, left + Inches(0.14), top + Inches(0.06),
                width - Inches(0.28), Inches(0.18),
                title, font=FONT_SANS, size_pt=9.5, bold=True,
                color=C_FG, name=(name or "card-hd") + "-title")
    if subtitle:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.22),
                    width - Inches(0.28), Inches(0.13),
                    subtitle, font=FONT_SANS, size_pt=8, color=C_MUTED,
                    name=(name or "card-hd") + "-sub")
    return top + h


# ══════════════════════════════════════════════════════════════════════════════
#  Table row helper
# ══════════════════════════════════════════════════════════════════════════════

def add_table_header_row(slide, left, top, width, cols, name="th"):
    row_h = Inches(0.24)
    add_rect(slide, left, top, width, row_h,
             fill_color=C_SURFACE, line_color=C_BORDER,
             line_width_pt=0.5, name=f"{name}-bg")
    x = left
    for col_label, col_w_frac in cols:
        col_w = width * col_w_frac
        add_textbox(slide, x + Inches(0.06), top + Inches(0.05),
                    col_w - Inches(0.08), Inches(0.14),
                    col_label.upper(), font=FONT_SANS, size_pt=7,
                    bold=True, color=C_MUTED, all_caps=True,
                    letter_spacing_pt=1.0, name=f"{name}-{col_label}")
        x += col_w
    return top + row_h

def add_table_data_row(slide, left, top, width, cells, row_h=Inches(0.36),
                       bg_color=None, name="tr"):
    add_rect(slide, left, top, width, row_h,
             fill_color=bg_color or C_SURFACE,
             line_color=C_BG, line_width_pt=0.3, name=f"{name}-bg")
    x = left
    for i, (text, col_w_frac, kwargs) in enumerate(cells):
        col_w = width * col_w_frac
        add_textbox(slide, x + Inches(0.06), top + Inches(0.05),
                    col_w - Inches(0.08), row_h - Inches(0.10),
                    text, font=kwargs.get("font", FONT_SANS),
                    size_pt=kwargs.get("size_pt", 9),
                    bold=kwargs.get("bold", False),
                    color=kwargs.get("color", C_FG),
                    name=f"{name}-cell-{i}")
        x += col_w
    return top + row_h


# ══════════════════════════════════════════════════════════════════════════════
#  Progress bar
# ══════════════════════════════════════════════════════════════════════════════

def add_progress_bar(slide, left, top, width, fill_pct, fill_color=None, name="prog"):
    add_rect(slide, left, top, width, Inches(0.055),
             fill_color=C_BORDER, name=f"{name}-track")
    fill_w = int(width * fill_pct / 100)
    if fill_w > 10:
        add_rect(slide, left, top, fill_w, Inches(0.055),
                 fill_color=fill_color or C_NAVY, name=f"{name}-fill")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_NAVY)
    # Top accent stripe
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.055), fill_color=C_ACCENT,
             name="cover-stripe")

    BLOCKS = [
        (Inches(0.20), Inches(0.22)),  # eyebrow
        (Inches(0.58), Inches(0.18)),  # title
        (Inches(0.28), Inches(0.26)),  # subtitle
        (Inches(0.18), Inches(0.34)),  # desc
        (Inches(0.16), Inches(0.00)),  # meta
    ]
    c = hero_layout(BLOCKS)

    cx = CANVAS_W / 2 - Inches(4.5)
    W = Inches(9)

    t = c.take(*BLOCKS[0], label="eyebrow")
    add_textbox(slide, cx, t, W, Inches(0.20),
                "ADMIN CONSOLE  ·  URL REDIRECT MANAGEMENT",
                font=FONT_MONO, size_pt=9, color=C_ACCENT,
                align=PP_ALIGN.CENTER, all_caps=False,
                letter_spacing_pt=2.0, name="cover-eyebrow")

    t = c.take(*BLOCKS[1], label="title")
    add_textbox(slide, cx, t, W, Inches(0.58),
                "PVM URL Admin",
                font=FONT_SANS, size_pt=48, bold=True,
                color=C_SURFACE, align=PP_ALIGN.CENTER,
                name="cover-title")

    t = c.take(*BLOCKS[2], label="subtitle")
    add_textbox(slide, cx, t, W, Inches(0.28),
                "Managed Redirects · go.pvm.co.za",
                font=FONT_MONO, size_pt=14, color=C_ACCENT,
                align=PP_ALIGN.CENTER, name="cover-subtitle")

    t = c.take(*BLOCKS[3], label="desc")
    add_textbox(slide, cx, t, W, Inches(0.28),
                "High-fidelity admin prototype — 5 screens",
                font=FONT_SANS, size_pt=12,
                color=RGBColor(0x9B,0xA0,0xA6),
                align=PP_ALIGN.CENTER, name="cover-desc")

    t = c.take(*BLOCKS[4], label="meta")
    add_textbox(slide, cx, t, W, Inches(0.16),
                "pvm.co.za  ·  2026-05-13",
                font=FONT_MONO, size_pt=8.5,
                color=RGBColor(0x72,0x80,0x90),
                align=PP_ALIGN.CENTER, name="cover-meta")

    add_footer(slide, "PVM URL Admin", "1 / 6", theme="dark")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Dashboard
# ══════════════════════════════════════════════════════════════════════════════

def slide_dashboard(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_BG)
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.04), fill_color=C_ACCENT,
             name="stripe")
    add_sidebar(slide, active_view="dashboard")
    add_topbar(slide, "PVM URL Admin  ›  Dashboard", 2, 6)

    CX = SIDEBAR_W + Inches(0.26)
    CW = CANVAS_W - CX - Inches(0.22)
    c = Cursor(y_start=Inches(0.55), cap=CONTENT_MAX_Y)

    add_page_header(slide, "Admin Console", "Dashboard",
                    "Overview of redirects and recent activity.", c)

    # ── KPI row ───────────────────────────────────────────────────────────────
    kpi_y = c.take(Inches(0.88), gap=Inches(0.14), label="kpi-row")
    kpi_w = (CW - Inches(0.36)) / 4
    kpis = [
        ("Total Redirects", "42", "3 added this week", C_GREEN),
        ("Total Clicks",    "1,287", "127 this week", C_GREEN),
        ("Clicks Today",    "14", "As of 10:32 AM", C_MUTED),
        ("Active",          "38", "4 paused or expired", C_MUTED),
    ]
    for i, (lbl, val, delta, dcol) in enumerate(kpis):
        kx = CX + i * (kpi_w + Inches(0.12))
        add_kpi_card(slide, kx, kpi_y, kpi_w, Inches(0.88),
                     lbl, val, delta, dcol)

    # ── Two-col row ───────────────────────────────────────────────────────────
    col_y = c.take(Inches(2.52), gap=Inches(0.14), label="two-col")
    col_w = (CW - Inches(0.14)) / 2
    COL_H = Inches(2.52)

    # Top Redirects card
    left_x = CX
    add_card_bg(slide, left_x, col_y, col_w, COL_H, name="top-redirects-card")
    body_top = add_card_header(slide, left_x, col_y, col_w,
                               "Top Redirects", "By total clicks, all time",
                               name="top-redirects-hd")
    rows = [
        ("PVM Energy Bar", "go.pvm.co.za/4frr2zc5", "Campaign", "286"),
        ("Octane Pre-workout", "go.pvm.co.za/oct-22", "Print / QR", "214"),
        ("Protein XTR — Combo", "go.pvm.co.za/xtr-cmb", "Referral", "198"),
        ("Fusion Meal Replacement", "go.pvm.co.za/ptn-1", "Print / QR", "147"),
        ("Reignite — Recovery Stack", "go.pvm.co.za/rgn-r", "Event", "112"),
    ]
    ry = body_top + Inches(0.06)
    row_h = Inches(0.40)
    for i, (title, code, purpose, clicks) in enumerate(rows):
        add_rect(slide, left_x, ry, col_w, row_h,
                 fill_color=C_SURFACE, line_color=C_BG, line_width_pt=0.3,
                 name=f"top-row-{i}-bg")
        add_textbox(slide, left_x + Inches(0.12), ry + Inches(0.04),
                    col_w * 0.55, Inches(0.15), title,
                    font=FONT_SANS, size_pt=8.5, bold=True, color=C_FG,
                    name=f"top-row-{i}-title")
        add_textbox(slide, left_x + Inches(0.12), ry + Inches(0.20),
                    col_w * 0.55, Inches(0.13), code,
                    font=FONT_MONO, size_pt=7, color=C_TEAL,
                    name=f"top-row-{i}-code")
        add_badge(slide, left_x + col_w * 0.60, ry + Inches(0.11),
                  purpose, name=f"top-row-{i}-badge")
        add_textbox(slide, left_x + col_w * 0.87, ry + Inches(0.09),
                    col_w * 0.12, Inches(0.18), clicks,
                    font=FONT_SANS, size_pt=9.5, bold=True, color=C_FG,
                    align=PP_ALIGN.RIGHT, name=f"top-row-{i}-clicks")
        ry += row_h

    # Recent Activity card
    right_x = CX + col_w + Inches(0.14)
    add_card_bg(slide, right_x, col_y, col_w, COL_H, name="recent-card")
    act_body_top = add_card_header(slide, right_x, col_y, col_w,
                                   "Recent Activity", "Latest clicks and changes",
                                   name="recent-hd")
    activities = [
        ("PVM Energy Bar — matched click", "go.pvm.co.za/4frr2zc5 · No referrer", "10:32", C_GREEN),
        ("New redirect — Reactor Sample Pack", "go.pvm.co.za/rct-sp · Print / QR", "09:15", C_TEAL),
        ("Octane Pre-workout — matched click", "go.pvm.co.za/oct-22 · instagram.com", "08:47", C_GREEN),
        ("Fusion Meal Replacement — dest updated", "Edited by Marius", "Yesterday", C_AMBER),
        ("Protein XTR Combo — matched click", "go.pvm.co.za/xtr-cmb · google.com", "Yesterday", C_GREEN),
        ("Summer Promo 2025 — redirect paused", "Campaign expired", "2 days ago", C_MUTED),
    ]
    ay = act_body_top + Inches(0.06)
    act_h = Inches(0.36)
    for i, (main, meta, time, dot_col) in enumerate(activities):
        add_rect(slide, right_x + Inches(0.14), ay + Inches(0.12),
                 Inches(0.08), Inches(0.08),
                 fill_color=dot_col, name=f"act-dot-{i}")
        add_textbox(slide, right_x + Inches(0.30), ay + Inches(0.04),
                    col_w - Inches(0.80), Inches(0.16), main,
                    font=FONT_SANS, size_pt=8, bold=False, color=C_FG,
                    name=f"act-main-{i}")
        add_textbox(slide, right_x + Inches(0.30), ay + Inches(0.20),
                    col_w - Inches(0.80), Inches(0.12), meta,
                    font=FONT_MONO, size_pt=7, color=C_MUTED,
                    name=f"act-meta-{i}")
        add_textbox(slide, right_x + col_w - Inches(0.90), ay + Inches(0.08),
                    Inches(0.76), Inches(0.14), time,
                    font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    align=PP_ALIGN.RIGHT, name=f"act-time-{i}")
        ay += act_h

    # ── Purpose breakdown card ─────────────────────────────────────────────────
    pb_y = c.take(Inches(1.00), gap=0, label="purpose-card")
    add_card_bg(slide, CX, pb_y, CW, Inches(1.00), name="purpose-card")
    add_card_header(slide, CX, pb_y, CW,
                    "Purpose breakdown", "Distribution across 42 redirects",
                    name="purpose-hd")
    purposes = [
        ("Print / QR", 43, "18 redirects", C_TEAL),
        ("Campaign",   29, "12 redirects", C_AMBER),
        ("Referral",   19, "8 redirects",  C_GREEN),
        ("Event",      10, "4 redirects",  C_PURPLE),
    ]
    pb_body_y = pb_y + Inches(0.38)
    pcol_w = (CW - Inches(0.42)) / 4
    for i, (lbl, pct, count, col) in enumerate(purposes):
        px = CX + Inches(0.14) + i * (pcol_w + Inches(0.14))
        add_badge(slide, px, pb_body_y, lbl, name=f"pb-badge-{i}")
        add_textbox(slide, px + Inches(0.78), pb_body_y + Inches(0.02),
                    pcol_w - Inches(0.82), Inches(0.16),
                    str(pct if lbl == "Campaign" else (18 if lbl == "Print / QR" else (8 if lbl=="Referral" else 4))),
                    font=FONT_SANS, size_pt=9, bold=True, color=C_FG,
                    align=PP_ALIGN.RIGHT, name=f"pb-count-{i}")
        add_progress_bar(slide, px, pb_body_y + Inches(0.24),
                         pcol_w - Inches(0.02), pct, fill_color=col,
                         name=f"pb-prog-{i}")
        add_textbox(slide, px, pb_body_y + Inches(0.36),
                    pcol_w - Inches(0.02), Inches(0.14),
                    f"{pct}% · {count}",
                    font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    name=f"pb-label-{i}")

    add_footer(slide, "Screen 01 — Dashboard", "2 / 6")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — All Redirects
# ══════════════════════════════════════════════════════════════════════════════

def slide_redirects(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_BG)
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.04), fill_color=C_ACCENT, name="stripe")
    add_sidebar(slide, active_view="redirects")
    add_topbar(slide, "PVM URL Admin  ›  All Redirects", 3, 6)

    CX = SIDEBAR_W + Inches(0.26)
    CW = CANVAS_W - CX - Inches(0.22)
    c = Cursor(y_start=Inches(0.55), cap=CONTENT_MAX_Y)
    add_page_header(slide, "Admin Console", "Redirects",
                    "Manage stable printed, QR, and campaign URLs.", c)

    card_top = c.take(Inches(4.80), gap=0, label="redirects-card")
    add_card_bg(slide, CX, card_top, CW, Inches(4.80), name="redirects-card")

    # Filter bar
    fb_y = card_top + Inches(0.08)
    add_rect(slide, CX + Inches(0.12), fb_y,
             Inches(2.8), Inches(0.28),
             fill_color=C_SURFACE, line_color=C_BORDER, line_width_pt=0.5,
             name="search-box")
    add_textbox(slide, CX + Inches(0.26), fb_y + Inches(0.06),
                Inches(2.5), Inches(0.16),
                "Code, title, tag, or destination…",
                font=FONT_SANS, size_pt=8.5, color=C_MUTED, name="search-placeholder")
    for i, (lbl, w) in enumerate([("All categories", 1.14), ("All purposes", 1.06), ("All statuses", 1.02)]):
        add_rect(slide, CX + Inches(3.10 + i * 1.22), fb_y,
                 Inches(w), Inches(0.28),
                 fill_color=C_SURFACE, line_color=C_BORDER, line_width_pt=0.5,
                 name=f"filter-{i}")
        add_textbox(slide, CX + Inches(3.16 + i * 1.22), fb_y + Inches(0.06),
                    Inches(w - 0.12), Inches(0.16), lbl,
                    font=FONT_SANS, size_pt=8, color=C_FG, name=f"filter-lbl-{i}")

    # Table
    th_y = card_top + Inches(0.44)
    COL_DEFS = [
        ("Short URL", 0.22), ("Title / Destination", 0.26), ("Category", 0.10),
        ("Purpose", 0.10), ("Tags", 0.12), ("Clicks", 0.08), ("Status", 0.07), ("Updated", 0.05),
    ]
    add_table_header_row(slide, CX, th_y, CW, COL_DEFS, name="th")

    redirects = [
        ("go.pvm.co.za/4frr2zc5", "PVM Energy Bar", "General", "Campaign", "packaging, social", "286", "Active", "2026/05/13"),
        ("go.pvm.co.za/oct-22",   "Octane Pre-workout", "Products", "Print / QR", "packaging", "214", "Active", "2026/05/12"),
        ("go.pvm.co.za/xtr-cmb",  "Protein XTR — Combo", "Products", "Referral", "influencer", "198", "Active", "2026/05/10"),
        ("go.pvm.co.za/ptn-1",    "Fusion Meal Replacement", "General", "Print / QR", "packaging", "147", "Active", "2026/05/13"),
        ("go.pvm.co.za/rgn-r",    "Reignite — Recovery Stack", "Events", "Event", "race-2025", "112", "Active", "2026/05/09"),
        ("go.pvm.co.za/rct-sp",   "Reactor Sample Pack", "Products", "Print / QR", "packaging", "67", "Active", "2026/05/13"),
        ("go.pvm.co.za/smr-25",   "Summer Promo 2025", "Campaigns", "Campaign", "social", "89", "Expired", "2026/04/30"),
        ("go.pvm.co.za/tri-ref",  "Triathlon Ref — @runwithme", "Partners", "Referral", "influencer", "44", "Paused", "2026/05/01"),
    ]

    ry = th_y + Inches(0.24)
    row_h = Inches(0.42)
    for i, (code, title, cat, purpose, tags, clicks, status, updated) in enumerate(redirects):
        if ry + row_h > card_top + Inches(4.70):
            break
        bg = C_BG if i % 2 == 1 else C_SURFACE
        add_rect(slide, CX, ry, CW, row_h,
                 fill_color=bg, line_color=C_BG, line_width_pt=0.2,
                 name=f"tr-{i}-bg")
        cols = [
            (code,    0.22, {"font": FONT_MONO, "size_pt": 7.5, "color": C_TEAL}),
            (title,   0.26, {"size_pt": 8.5, "bold": True}),
            (cat,     0.10, {"size_pt": 8, "color": C_MUTED}),
            (purpose, 0.10, {"size_pt": 8}),
            (tags,    0.12, {"size_pt": 7.5, "color": C_MUTED}),
            (clicks,  0.08, {"size_pt": 9, "bold": True}),
            (status,  0.07, {"size_pt": 8}),
            (updated, 0.05, {"size_pt": 7.5, "color": C_MUTED}),
        ]
        x = CX
        for j, (txt, frac, kw) in enumerate(cols):
            cw = CW * frac
            add_textbox(slide, x + Inches(0.06), ry + Inches(0.10),
                        cw - Inches(0.08), row_h - Inches(0.14),
                        txt, font=kw.get("font", FONT_SANS),
                        size_pt=kw.get("size_pt", 8.5),
                        bold=kw.get("bold", False),
                        color=kw.get("color", C_FG),
                        name=f"tr-{i}-col-{j}")
            x += cw
        ry += row_h

    # Pagination
    pg_y = card_top + Inches(4.52)
    add_textbox(slide, CX + Inches(0.14), pg_y, Inches(2), Inches(0.18),
                "Showing 8 of 42 redirects",
                font=FONT_SANS, size_pt=8, color=C_MUTED, name="pg-count")
    for i, lbl in enumerate(["← Prev", "1", "2", "3", "Next →"]):
        btn_w = Inches(0.44) if len(lbl) > 2 else Inches(0.28)
        bx = CANVAS_W - Inches(0.28) - (5-i) * Inches(0.48)
        bg = C_NAVY if lbl == "1" else C_SURFACE
        fc = C_SURFACE if lbl == "1" else C_FG
        add_rect(slide, bx, pg_y - Inches(0.02), btn_w, Inches(0.24),
                 fill_color=bg, line_color=C_BORDER, line_width_pt=0.4,
                 name=f"pg-btn-{i}")
        add_textbox(slide, bx, pg_y - Inches(0.01), btn_w, Inches(0.22),
                    lbl, font=FONT_SANS, size_pt=8, bold=(lbl=="1"),
                    color=fc, align=PP_ALIGN.CENTER, name=f"pg-btn-lbl-{i}")

    add_footer(slide, "Screen 02 — All Redirects", "3 / 6")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Create / Edit Redirect
# ══════════════════════════════════════════════════════════════════════════════

def slide_create(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_BG)
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.04), fill_color=C_ACCENT, name="stripe")
    add_sidebar(slide, active_view="create")
    add_topbar(slide, "PVM URL Admin  ›  Redirects  ›  New Redirect", 4, 6)

    CX = SIDEBAR_W + Inches(0.26)
    CW = CANVAS_W - CX - Inches(0.22)
    c = Cursor(y_start=Inches(0.55), cap=CONTENT_MAX_Y)
    add_page_header(slide, "Redirects", "New redirect",
                    "Create a stable short URL. The code is permanent — the destination can always be updated.", c)

    # Layout: main form (left) + sidebar panels (right)
    FORM_W = CW * 0.68
    SIDE_W = CW * 0.30
    SIDE_X = CX + FORM_W + CW * 0.02

    form_top = c.y

    # Card 1: Basic information
    c1_h = Inches(1.74)
    add_card_bg(slide, CX, form_top, FORM_W, c1_h, name="card-basic")
    add_card_header(slide, CX, form_top, FORM_W, "Basic information", name="card-basic-hd")
    fields = [
        ("Short code — permanent once created", "go.pvm.co.za/ [e.g. oct-22]", Inches(0.38)),
        ("Title", "e.g. Fusion Meal Replacement — Packaging Q1 2025", Inches(0.36)),
        ("Destination URL", "https://pvm.co.za/product/…", Inches(0.36)),
    ]
    fy = form_top + Inches(0.36)
    for lbl, placeholder, fh in fields:
        add_textbox(slide, CX + Inches(0.14), fy, FORM_W - Inches(0.28),
                    Inches(0.16), lbl, font=FONT_SANS, size_pt=8.5,
                    bold=True, color=C_FG, name=f"field-lbl-{lbl[:10]}")
        fy += Inches(0.18)
        add_rect(slide, CX + Inches(0.14), fy, FORM_W - Inches(0.28),
                 Inches(0.26), fill_color=C_SURFACE, line_color=C_BORDER,
                 line_width_pt=0.5, name=f"field-input-{lbl[:10]}")
        add_textbox(slide, CX + Inches(0.22), fy + Inches(0.05),
                    FORM_W - Inches(0.44), Inches(0.16), placeholder,
                    font=FONT_MONO if "go.pvm" in placeholder else FONT_SANS,
                    size_pt=8, color=C_MUTED, name=f"field-ph-{lbl[:10]}")
        fy += fh

    # Card 2: Classification
    c2_top = form_top + c1_h + Inches(0.12)
    c2_h = Inches(1.62)
    add_card_bg(slide, CX, c2_top, FORM_W, c2_h, name="card-class")
    add_card_header(slide, CX, c2_top, FORM_W, "Classification",
                    "Category, purpose, and tags", name="card-class-hd")
    # Category + Purpose selects (2-up)
    hy = c2_top + Inches(0.44)
    for i, (lbl, opts) in enumerate([
        ("Category", "Select…"),
        ("Purpose", "Select…"),
    ]):
        hx = CX + Inches(0.14) + i * (FORM_W / 2 - Inches(0.07))
        hw = FORM_W / 2 - Inches(0.21)
        add_textbox(slide, hx, hy, hw, Inches(0.16), lbl,
                    font=FONT_SANS, size_pt=8.5, bold=True, color=C_FG,
                    name=f"class-lbl-{i}")
        add_rect(slide, hx, hy + Inches(0.18), hw, Inches(0.26),
                 fill_color=C_SURFACE, line_color=C_BORDER, line_width_pt=0.5,
                 name=f"class-sel-{i}")
        add_textbox(slide, hx + Inches(0.08), hy + Inches(0.23),
                    hw - Inches(0.16), Inches(0.16), opts,
                    font=FONT_SANS, size_pt=8.5, color=C_MUTED,
                    name=f"class-sel-ph-{i}")
    # Tags selector
    ty = hy + Inches(0.54)
    add_textbox(slide, CX + Inches(0.14), ty, FORM_W - Inches(0.28),
                Inches(0.16), "Tags — select all that apply",
                font=FONT_SANS, size_pt=8.5, bold=True, color=C_FG,
                name="tags-lbl")
    ty += Inches(0.18)
    tags = ["packaging", "product-page", "social", "influencer", "race-2025", "print-run-q1", "email", "instagram"]
    tx = CX + Inches(0.14)
    for i, tag in enumerate(tags):
        tw = Inches(0.82)
        bg = C_NAVY if tag == "packaging" else C_BG
        fc = C_SURFACE if tag == "packaging" else C_MUTED
        add_rect(slide, tx, ty, tw, Inches(0.22),
                 fill_color=bg, line_color=C_BORDER, line_width_pt=0.4,
                 name=f"tag-opt-{i}")
        add_textbox(slide, tx, ty, tw, Inches(0.22), tag,
                    font=FONT_SANS, size_pt=7.5, color=fc,
                    align=PP_ALIGN.CENTER, name=f"tag-opt-lbl-{i}")
        tx += tw + Inches(0.07)
        if tx + tw > CX + FORM_W - Inches(0.14):
            tx = CX + Inches(0.14)
            ty += Inches(0.28)

    # Right sidebar panels
    sy = form_top

    # Status panel
    add_card_bg(slide, SIDE_X, sy, SIDE_W, Inches(0.72), name="status-card")
    add_card_header(slide, SIDE_X, sy, SIDE_W, "Status", name="status-hd")
    add_rect(slide, SIDE_X + Inches(0.14), sy + Inches(0.44),
             Inches(0.34), Inches(0.18),
             fill_color=C_NAVY, name="toggle-track")
    add_rect(slide, SIDE_X + Inches(0.30), sy + Inches(0.46),
             Inches(0.12), Inches(0.14),
             fill_color=C_SURFACE, name="toggle-thumb")
    add_textbox(slide, SIDE_X + Inches(0.56), sy + Inches(0.42),
                SIDE_W - Inches(0.70), Inches(0.16),
                "Active", font=FONT_SANS, size_pt=9, bold=True,
                color=C_FG, name="status-lbl")
    add_textbox(slide, SIDE_X + Inches(0.56), sy + Inches(0.58),
                SIDE_W - Inches(0.70), Inches(0.13),
                "Redirect is live and resolving",
                font=FONT_SANS, size_pt=7.5, color=C_MUTED, name="status-sub")

    sy += Inches(0.84)

    # Short URL preview
    add_card_bg(slide, SIDE_X, sy, SIDE_W, Inches(0.82), name="preview-card")
    add_card_header(slide, SIDE_X, sy, SIDE_W, "Short URL preview", name="preview-hd")
    add_rect(slide, SIDE_X + Inches(0.14), sy + Inches(0.44),
             SIDE_W - Inches(0.28), Inches(0.26),
             fill_color=C_BG, line_color=C_BORDER, line_width_pt=0.5,
             name="preview-box")
    add_textbox(slide, SIDE_X + Inches(0.22), sy + Inches(0.49),
                SIDE_W - Inches(0.44), Inches(0.16),
                "go.pvm.co.za/  —",
                font=FONT_MONO, size_pt=8.5, color=C_TEAL, name="preview-txt")

    sy += Inches(0.96)

    # QR Code panel
    add_card_bg(slide, SIDE_X, sy, SIDE_W, Inches(1.30), name="qr-card")
    add_card_header(slide, SIDE_X, sy, SIDE_W, "QR Code", name="qr-hd")
    add_rect(slide, SIDE_X + SIDE_W/2 - Inches(0.64), sy + Inches(0.44),
             Inches(1.28), Inches(0.80),
             fill_color=C_BG, line_color=C_BORDER, line_width_pt=1.0,
             name="qr-stub")
    add_textbox(slide, SIDE_X + Inches(0.14), sy + Inches(0.58),
                SIDE_W - Inches(0.28), Inches(0.52),
                "QR generation arrives in the next release.",
                font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                align=PP_ALIGN.CENTER, name="qr-note")
    # Coming soon badge
    add_rect(slide, SIDE_X + SIDE_W - Inches(0.90), sy + Inches(0.08),
             Inches(0.82), Inches(0.17),
             fill_color=C_BG, line_color=C_BORDER, line_width_pt=0.4,
             name="coming-soon-bg")
    add_textbox(slide, SIDE_X + SIDE_W - Inches(0.90), sy + Inches(0.08),
                Inches(0.82), Inches(0.17),
                "Coming soon", font=FONT_SANS, size_pt=7, color=C_MUTED,
                align=PP_ALIGN.CENTER, name="coming-soon-lbl")

    add_footer(slide, "Screen 03 — Create / Edit Redirect", "4 / 6")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Redirect Detail / Analytics
# ══════════════════════════════════════════════════════════════════════════════

def slide_detail(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_BG)
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.04), fill_color=C_ACCENT, name="stripe")
    add_sidebar(slide, active_view="redirects")
    add_topbar(slide, "PVM URL Admin  ›  Redirects  ›  PVM Energy Bar", 5, 6)

    CX = SIDEBAR_W + Inches(0.26)
    CW = CANVAS_W - CX - Inches(0.22)
    c = Cursor(y_start=Inches(0.58), cap=CONTENT_MAX_Y)

    # Back link
    t = c.take(Inches(0.18), gap=Inches(0.10), label="back")
    add_textbox(slide, CX, t, Inches(1.5), Inches(0.18),
                "← Back to redirects", font=FONT_SANS, size_pt=8.5,
                color=C_MUTED, name="back-link")

    # Page title section
    t = c.take(Inches(0.18), gap=Inches(0.04), label="badges")
    add_badge(slide, CX, t, "Campaign", name="detail-badge-campaign")
    add_badge(slide, CX + Inches(0.80), t, "Active", name="detail-badge-active")
    t = c.take(Inches(0.32), gap=Inches(0.06), label="title")
    add_textbox(slide, CX, t, CW * 0.75, Inches(0.32),
                "PVM Energy Bar", font=FONT_SANS, size_pt=26,
                bold=True, color=C_FG, name="detail-title")
    t = c.take(Inches(0.22), gap=Inches(0.06), label="url-row")
    add_rect(slide, CX, t, Inches(2.60), Inches(0.26),
             fill_color=C_BG, line_color=C_BORDER, line_width_pt=0.5,
             name="detail-url-box")
    add_textbox(slide, CX + Inches(0.10), t + Inches(0.05),
                Inches(2.40), Inches(0.16),
                "go.pvm.co.za/4frr2zc5",
                font=FONT_MONO, size_pt=8.5, color=C_TEAL, name="detail-url")
    add_textbox(slide, CX, t + Inches(0.28), CW * 0.75, Inches(0.18),
                "→  https://pvm.co.za/product/pvm-energy-bar/",
                font=FONT_MONO, size_pt=8, color=C_MUTED, name="detail-dest")
    c.take(Inches(0.18), gap=Inches(0.12), label="dest-row")

    # KPI row
    kpi_y = c.take(Inches(0.88), gap=Inches(0.14), label="detail-kpis")
    kpi_w = (CW - Inches(0.36)) / 4
    detail_kpis = [
        ("Total Clicks", "286", "+23 this week", C_GREEN),
        ("Last Click", "Today\n10:32 AM", "", C_MUTED),
        ("Created", "2026/05/13", "By Marius", C_MUTED),
        ("Code", "4frr2zc5", "Stable · Print / QR", C_MUTED),
    ]
    for i, (lbl, val, delta, dc) in enumerate(detail_kpis):
        kx = CX + i * (kpi_w + Inches(0.12))
        add_kpi_card(slide, kx, kpi_y, kpi_w, Inches(0.88), lbl, val, delta, dc)

    # Two-col: chart + referrers
    row2_y = c.take(Inches(1.90), gap=Inches(0.12), label="analytics-row")
    col_w = (CW - Inches(0.14)) / 2

    # Click activity card (left)
    add_card_bg(slide, CX, row2_y, col_w, Inches(1.90), name="chart-card")
    add_card_header(slide, CX, row2_y, col_w,
                    "Click activity", "Last 14 days · 286 total", name="chart-hd")
    # SVG-style sparkline using thin rectangles
    chart_x = CX + Inches(0.14)
    chart_y = row2_y + Inches(0.44)
    chart_w = col_w - Inches(0.28)
    chart_h = Inches(1.20)
    add_rect(slide, chart_x, chart_y, chart_w, chart_h,
             fill_color=RGBColor(0xF8,0xFA,0xFC), line_color=C_BORDER,
             line_width_pt=0.3, name="chart-bg")
    # Sparkline bars
    bar_data = [88, 79, 92, 71, 60, 65, 43, 52, 37, 25, 15, 6, 4, 4]
    bw = chart_w / (len(bar_data) * 1.6)
    max_val = max(bar_data)
    for i, v in enumerate(bar_data):
        bh = chart_h * 0.85 * v / max_val
        bx = chart_x + Inches(0.06) + i * (bw + Inches(0.04))
        by = chart_y + chart_h - bh - Inches(0.04)
        alpha = 0.4 + 0.6 * v / max_val
        col_val = int(13 + (1-alpha) * 220)
        add_rect(slide, bx, by, bw, bh,
                 fill_color=RGBColor(col_val, 31+int((1-alpha)*60), 53+int((1-alpha)*100)),
                 name=f"bar-{i}")

    # Referrers card (right)
    ref_x = CX + col_w + Inches(0.14)
    add_card_bg(slide, ref_x, row2_y, col_w, Inches(1.90), name="ref-card")
    add_card_header(slide, ref_x, row2_y, col_w,
                    "Top referrers", "Source of 286 clicks", name="ref-hd")
    referrers = [
        ("Direct / No referrer", "120 · 42%", 42, C_NAVY),
        ("instagram.com",        "80 · 28%",  28, C_TEAL),
        ("google.com",           "40 · 14%",  14, C_NAVY),
        ("facebook.com",         "26 · 9%",    9, C_NAVY),
        ("email / newsletter",   "20 · 7%",    7, C_TEAL),
    ]
    rfy = row2_y + Inches(0.44)
    for src, cnt, pct, col in referrers:
        add_textbox(slide, ref_x + Inches(0.14), rfy, col_w * 0.58,
                    Inches(0.16), src, font=FONT_SANS, size_pt=8.5,
                    bold=True, color=C_FG, name=f"ref-src-{src[:6]}")
        add_textbox(slide, ref_x + col_w * 0.60, rfy, col_w * 0.36,
                    Inches(0.16), cnt, font=FONT_SANS, size_pt=8,
                    color=C_MUTED, align=PP_ALIGN.RIGHT, name=f"ref-cnt-{src[:6]}")
        add_progress_bar(slide, ref_x + Inches(0.14), rfy + Inches(0.18),
                         col_w - Inches(0.28), pct, fill_color=col,
                         name=f"ref-prog-{src[:6]}")
        rfy += Inches(0.34)

    # Recent clicks table
    tbl_y = c.take(Inches(1.02), gap=0, label="clicks-table")
    add_card_bg(slide, CX, tbl_y, CW, Inches(1.02), name="clicks-card")
    add_card_header(slide, CX, tbl_y, CW,
                    "Recent clicks", "Last 20 recorded events", name="clicks-hd")
    cl_cols = [("Referrer", 0.28), ("Device / Browser", 0.24),
               ("Result", 0.18), ("Timestamp", 0.30)]
    cth_y = add_table_header_row(slide, CX, tbl_y + Inches(0.36),
                                 CW, cl_cols, name="clicks-th")
    click_rows = [
        ("No referrer", "Chrome · Win 10", "Matched", "2026/05/13 10:32:19"),
        ("instagram.com", "Safari · iOS 18", "Matched", "2026/05/13 09:14:07"),
        ("google.com", "Chrome · Android", "Matched", "2026/05/12 16:43:52"),
    ]
    cly = cth_y
    for i, (ref, device, result, ts) in enumerate(click_rows):
        if cly + Inches(0.24) > tbl_y + Inches(1.00):
            break
        x = CX
        for txt, frac in [(ref, 0.28), (device, 0.24), (result, 0.18), (ts, 0.30)]:
            cw2 = CW * frac
            col = C_TEAL if txt == "Matched" else (C_MUTED if txt == ts else C_FG)
            font = FONT_MONO if (txt == ref or txt == ts) else FONT_SANS
            add_textbox(slide, x + Inches(0.08), cly + Inches(0.04),
                        cw2 - Inches(0.10), Inches(0.18),
                        txt, font=font, size_pt=7.5, color=col,
                        name=f"cr-{i}-{txt[:6]}")
            x += cw2
        add_rect(slide, CX, cly + Inches(0.22), CW, Inches(0.01),
                 fill_color=C_BG, name=f"cr-div-{i}")
        cly += Inches(0.23)

    add_footer(slide, "Screen 04 — Redirect Detail / Analytics", "5 / 6")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Tags & Categories
# ══════════════════════════════════════════════════════════════════════════════

def slide_tags(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    paint_bg(slide, C_BG)
    add_rect(slide, 0, 0, CANVAS_W, Inches(0.04), fill_color=C_ACCENT, name="stripe")
    add_sidebar(slide, active_view="tags")
    add_topbar(slide, "PVM URL Admin  ›  Organisation  ›  Tags & Categories", 6, 6)

    CX = SIDEBAR_W + Inches(0.26)
    CW = CANVAS_W - CX - Inches(0.22)
    c = Cursor(y_start=Inches(0.55), cap=CONTENT_MAX_Y)
    add_page_header(slide, "Organisation", "Tags & Categories",
                    "Organise redirects for easier filtering and reporting.", c)

    col_top = c.take(Inches(2.54), gap=Inches(0.12), label="cats-tags-row")
    col_w = (CW - Inches(0.14)) / 2
    COL_H = Inches(2.54)

    # Categories card (left)
    add_card_bg(slide, CX, col_top, col_w, COL_H, name="cats-card")
    body_y = add_card_header(slide, CX, col_top, col_w,
                             "Categories", "High-level groupings · 5 total",
                             name="cats-hd")
    # Add input row
    add_rect(slide, CX + Inches(0.12), body_y + Inches(0.08),
             col_w - Inches(0.84), Inches(0.26),
             fill_color=C_SURFACE, line_color=C_BORDER, line_width_pt=0.5,
             name="cat-input")
    add_textbox(slide, CX + Inches(0.20), body_y + Inches(0.14),
                col_w - Inches(1.0), Inches(0.14),
                "New category name…", font=FONT_SANS, size_pt=8,
                color=C_MUTED, name="cat-input-ph")
    add_rect(slide, CX + col_w - Inches(0.64), body_y + Inches(0.08),
             Inches(0.52), Inches(0.26),
             fill_color=C_NAVY, name="cat-add-btn")
    add_textbox(slide, CX + col_w - Inches(0.64), body_y + Inches(0.08),
                Inches(0.52), Inches(0.26),
                "Add", font=FONT_SANS, size_pt=8, bold=True,
                color=C_SURFACE, align=PP_ALIGN.CENTER, name="cat-add-lbl")

    cats = [
        ("General", "12 redirects · Uncategorised or miscellaneous"),
        ("Products", "18 redirects · Product pages and PDPs"),
        ("Campaigns", "6 redirects · Promo and marketing campaigns"),
        ("Events", "4 redirects · Race expos and trade shows"),
        ("Partners", "2 redirects · Partner and affiliate links"),
    ]
    iy = body_y + Inches(0.44)
    row_h = Inches(0.38)
    for i, (name_txt, meta_txt) in enumerate(cats):
        if iy + row_h > col_top + COL_H - Inches(0.04):
            break
        add_rect(slide, CX, iy, col_w, row_h,
                 fill_color=C_SURFACE, line_color=C_BG, line_width_pt=0.2,
                 name=f"cat-row-{i}-bg")
        add_textbox(slide, CX + Inches(0.14), iy + Inches(0.06),
                    col_w * 0.55, Inches(0.16), name_txt,
                    font=FONT_SANS, size_pt=9, bold=True,
                    color=C_FG, name=f"cat-name-{i}")
        add_textbox(slide, CX + Inches(0.14), iy + Inches(0.22),
                    col_w * 0.55, Inches(0.13), meta_txt,
                    font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    name=f"cat-meta-{i}")
        # Edit / Delete buttons
        add_rect(slide, CX + col_w - Inches(1.24), iy + Inches(0.10),
                 Inches(0.52), Inches(0.20),
                 fill_color=C_SURFACE, line_color=C_BORDER,
                 line_width_pt=0.4, name=f"cat-edit-{i}")
        add_textbox(slide, CX + col_w - Inches(1.24), iy + Inches(0.10),
                    Inches(0.52), Inches(0.20),
                    "Edit", font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    align=PP_ALIGN.CENTER, name=f"cat-edit-lbl-{i}")
        add_rect(slide, CX + col_w - Inches(0.68), iy + Inches(0.10),
                 Inches(0.56), Inches(0.20),
                 fill_color=RGBColor(0xFE,0xF2,0xF2),
                 line_color=RGBColor(0xFE,0xCA,0xCA),
                 line_width_pt=0.4, name=f"cat-del-{i}")
        add_textbox(slide, CX + col_w - Inches(0.68), iy + Inches(0.10),
                    Inches(0.56), Inches(0.20),
                    "Delete", font=FONT_SANS, size_pt=7.5, color=C_RED,
                    align=PP_ALIGN.CENTER, name=f"cat-del-lbl-{i}")
        iy += row_h

    # Tags card (right)
    right_x = CX + col_w + Inches(0.14)
    add_card_bg(slide, right_x, col_top, col_w, COL_H, name="tags-card")
    tags_body_y = add_card_header(slide, right_x, col_top, col_w,
                                  "Tags", "Fine-grained labels · 8 total",
                                  name="tags-hd")
    # Add input row
    add_rect(slide, right_x + Inches(0.12), tags_body_y + Inches(0.08),
             col_w - Inches(0.84), Inches(0.26),
             fill_color=C_SURFACE, line_color=C_BORDER, line_width_pt=0.5,
             name="tag-input")
    add_textbox(slide, right_x + Inches(0.20), tags_body_y + Inches(0.14),
                col_w - Inches(1.0), Inches(0.14),
                "New tag…", font=FONT_SANS, size_pt=8,
                color=C_MUTED, name="tag-input-ph")
    add_rect(slide, right_x + col_w - Inches(0.64), tags_body_y + Inches(0.08),
             Inches(0.52), Inches(0.26),
             fill_color=C_NAVY, name="tag-add-btn")
    add_textbox(slide, right_x + col_w - Inches(0.64), tags_body_y + Inches(0.08),
                Inches(0.52), Inches(0.26),
                "Add", font=FONT_SANS, size_pt=8, bold=True,
                color=C_SURFACE, align=PP_ALIGN.CENTER, name="tag-add-lbl")

    tags_data = [
        ("packaging", "14 redirects"), ("product-page", "8 redirects"),
        ("social", "7 redirects"), ("influencer", "5 redirects"),
        ("print-run-q1", "6 redirects"), ("race-2025", "4 redirects"),
        ("email", "3 redirects"), ("instagram", "5 redirects"),
    ]
    tgy = tags_body_y + Inches(0.44)
    for i, (tag_name, tag_count) in enumerate(tags_data):
        if tgy + Inches(0.28) > col_top + COL_H - Inches(0.04):
            break
        add_rect(slide, right_x, tgy, col_w, Inches(0.28),
                 fill_color=C_SURFACE, line_color=C_BG, line_width_pt=0.2,
                 name=f"tag-row-{i}-bg")
        # Tag chip
        tw = Inches(0.86)
        add_rect(slide, right_x + Inches(0.14), tgy + Inches(0.05),
                 tw, Inches(0.18),
                 fill_color=C_BG, line_color=C_BORDER, line_width_pt=0.4,
                 name=f"tag-chip-{i}")
        add_textbox(slide, right_x + Inches(0.14), tgy + Inches(0.05),
                    tw, Inches(0.18), tag_name,
                    font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    align=PP_ALIGN.CENTER, name=f"tag-chip-lbl-{i}")
        add_textbox(slide, right_x + Inches(1.10), tgy + Inches(0.06),
                    Inches(0.80), Inches(0.16), tag_count,
                    font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    name=f"tag-count-{i}")
        # Rename / Delete
        add_rect(slide, right_x + col_w - Inches(1.30), tgy + Inches(0.06),
                 Inches(0.58), Inches(0.18),
                 fill_color=C_SURFACE, line_color=C_BORDER,
                 line_width_pt=0.4, name=f"tag-rename-{i}")
        add_textbox(slide, right_x + col_w - Inches(1.30), tgy + Inches(0.06),
                    Inches(0.58), Inches(0.18),
                    "Rename", font=FONT_SANS, size_pt=7, color=C_MUTED,
                    align=PP_ALIGN.CENTER, name=f"tag-rename-lbl-{i}")
        add_rect(slide, right_x + col_w - Inches(0.68), tgy + Inches(0.06),
                 Inches(0.56), Inches(0.18),
                 fill_color=RGBColor(0xFE,0xF2,0xF2),
                 line_color=RGBColor(0xFE,0xCA,0xCA),
                 line_width_pt=0.4, name=f"tag-del-{i}")
        add_textbox(slide, right_x + col_w - Inches(0.68), tgy + Inches(0.06),
                    Inches(0.56), Inches(0.18),
                    "Delete", font=FONT_SANS, size_pt=7, color=C_RED,
                    align=PP_ALIGN.CENTER, name=f"tag-del-lbl-{i}")
        tgy += Inches(0.28)

    # Purpose types card
    pt_y = c.take(Inches(1.36), gap=0, label="purpose-types")
    add_card_bg(slide, CX, pt_y, CW, Inches(1.36), name="purpose-card")
    add_card_header(slide, CX, pt_y, CW,
                    "Purpose types", "System-defined — not editable",
                    name="purpose-hd")
    purposes = [
        ("Print / QR", "bg-blue", C_TEAL, "For packaging, printed materials, and QR codes. Stable code — destination updateable."),
        ("Campaign",   "bg-amber", C_AMBER, "Short-lived URLs for promotions, product launches, or marketing campaigns."),
        ("Referral",   "bg-green", C_GREEN, "Tracks traffic from a specific source — influencer, partner, or affiliate."),
        ("Event",      "bg-purple", C_PURPLE, "Race expos, trade shows, demo days. Scoped to a single event."),
    ]
    pcol_w = (CW - Inches(0.42)) / 4
    for i, (lbl, _, badge_col, desc) in enumerate(purposes):
        px = CX + Inches(0.14) + i * (pcol_w + Inches(0.14))
        py = pt_y + Inches(0.42)
        add_badge(slide, px, py, lbl, name=f"pt-badge-{i}")
        add_textbox(slide, px, py + Inches(0.26), pcol_w, Inches(0.64),
                    desc, font=FONT_SANS, size_pt=7.5, color=C_MUTED,
                    word_wrap=True, name=f"pt-desc-{i}")

    add_footer(slide, "Screen 05 — Tags & Categories", "6 / 6")


# ══════════════════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = CANVAS_W
    prs.slide_height = CANVAS_H

    blank_layout = prs.slide_layouts[6]  # blank

    slide_cover(prs, blank_layout)
    slide_dashboard(prs, blank_layout)
    slide_redirects(prs, blank_layout)
    slide_create(prs, blank_layout)
    slide_detail(prs, blank_layout)
    slide_tags(prs, blank_layout)

    out_path = r"C:\Users\Marius\AppData\Roaming\Open Design\namespaces\release-stable-win\data\projects\c5052140-9f9e-477a-b1a6-fa22774c750f\pvm-url-admin.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    main()
